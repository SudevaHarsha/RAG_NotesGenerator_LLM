from pptx import Presentation
from typing import List, Dict
import os

class PPTExtractor:
    """
    Extract slides, text, speaker notes, and optionally images from PPTX.
    """

    @staticmethod
    def extract(ppt_path: str, output_image_dir: str = None) -> List[Dict]:
        prs = Presentation(ppt_path)
        slides_data = []
        os.makedirs(output_image_dir, exist_ok=True) if output_image_dir else None

        for idx, slide in enumerate(prs.slides, start=1):
            # Slide title
            title_shapes = [sh for sh in slide.shapes if sh.has_text_frame]
            title = title_shapes[0].text.strip() if title_shapes else ""

            # Bullet text
            bullet_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    bullet_texts.append("\n".join([p.text for p in text_frame.paragraphs]))
            bullet_text = "\n".join(bullet_texts)

            # Speaker notes
            notes_slide = slide.notes_slide
            speaker_notes = notes_slide.notes_text_frame.text if notes_slide else ""

            # Image extraction
            image_path = None
            if output_image_dir:
                image_path = os.path.join(output_image_dir, f"slide_{idx}.png")
                # Render slide as image (requires python-pptx + additional libs)
                from pptx2pdf import convert_single  # optional helper
                # Here we can fallback to pdf render if needed
                # For now, keep as None
                image_path = None

            slides_data.append({
                "slide_number": idx,
                "title": title,
                "bullet_text": bullet_text,
                "speaker_notes": speaker_notes,
                "slide_image_path": image_path
            })

        return slides_data
